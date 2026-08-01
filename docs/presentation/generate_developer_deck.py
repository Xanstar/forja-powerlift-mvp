"""Genera la guía editable para desarrollo del MVP piloto de Forja.

La dirección visual de este deck se limita a la presentación. Deriva de la
interfaz actual de Forja (grafito, tiza, acero y óxido moderado) para dar
coherencia al relato técnico; NO aprueba un rebranding final de la aplicación.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("Forja-MVP-Developer-Brief.pptx")

W = Inches(13.333)
H = Inches(7.5)

GRAPHITE = RGBColor(20, 22, 23)
GRAPHITE_2 = RGBColor(29, 32, 34)
STEEL = RGBColor(70, 76, 79)
STEEL_LIGHT = RGBColor(116, 123, 126)
CHALK = RGBColor(239, 235, 224)
CHALK_MUTED = RGBColor(184, 181, 172)
OXIDE = RGBColor(190, 73, 45)
OXIDE_DARK = RGBColor(121, 48, 33)
SUCCESS = RGBColor(104, 159, 117)
WARNING = RGBColor(213, 163, 72)

FONT_DISPLAY = "Arial"
FONT_BODY = "Arial"


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=CHALK, bold=False,
             font=FONT_BODY, align=PP_ALIGN.LEFT, margin=0, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=17, gap=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    frame.margin_top = frame.margin_bottom = 0
    for index, (label, body, color) in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = label
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
        r = p.add_run()
        r.text = body
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.color.rgb = CHALK
    return box


def add_line(slide, x1, y1, x2, y2, color=STEEL, width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def base_slide(prs, title, kicker, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = GRAPHITE
    add_rect(slide, Inches(0), Inches(0), Inches(0.13), H, OXIDE)
    add_text(slide, kicker.upper(), Inches(0.65), Inches(0.32), Inches(7.5), Inches(0.25),
             size=9, color=OXIDE, bold=True)
    add_text(slide, title, Inches(0.65), Inches(0.63), Inches(11.9), Inches(0.55),
             size=26, bold=True, font=FONT_DISPLAY)
    add_line(slide, Inches(0.65), Inches(1.25), Inches(12.68), Inches(1.25), STEEL, 1)
    add_text(slide, "Guía para desarrollo · MVP piloto", Inches(0.65), Inches(7.15), Inches(3.5), Inches(0.18),
             size=8, color=STEEL_LIGHT)
    add_text(slide, f"{number:02d}", Inches(12.1), Inches(7.12), Inches(0.55), Inches(0.22),
             size=9, color=CHALK_MUTED, bold=True, align=PP_ALIGN.RIGHT)
    return slide


def label_box(slide, title, body, x, y, w, h, accent=OXIDE, size=15):
    add_rect(slide, x, y, w, h, GRAPHITE_2, STEEL, radius=True)
    add_rect(slide, x, y, Inches(0.06), h, accent, accent)
    add_text(slide, title.upper(), x + Inches(0.22), y + Inches(0.16), w - Inches(0.38), Inches(0.22),
             size=9, color=accent, bold=True)
    add_text(slide, body, x + Inches(0.22), y + Inches(0.52), w - Inches(0.4), h - Inches(0.65),
             size=size, color=CHALK, bold=False)


def pill(slide, text, x, y, w, color=STEEL):
    add_rect(slide, x, y, w, Inches(0.34), color, color, radius=True)
    add_text(slide, text, x, y + Inches(0.01), w, Inches(0.28), size=9, color=CHALK,
             bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "Guía para desarrollo del MVP de Forja"
    prs.core_properties.subject = "Arquitectura, límite, ejecución y verificación del MVP piloto"
    prs.core_properties.author = "Equipo de desarrollo de Forja"

    # 1 — portada
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAPHITE
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, OXIDE)
    add_text(slide, "FORJA / MVP PILOTO", Inches(0.8), Inches(0.7), Inches(4), Inches(0.3),
             size=11, color=OXIDE, bold=True)
    add_text(slide, "Resumen de ejecución", Inches(0.8), Inches(1.25), Inches(7.6), Inches(0.9),
             size=46, bold=True, font=FONT_DISPLAY)
    add_text(slide, "Hacer confiable el ciclo de coaching antes de ampliarlo.",
             Inches(0.84), Inches(2.3), Inches(7.5), Inches(0.8), size=22, color=CHALK_MUTED)
    add_line(slide, Inches(0.82), Inches(3.35), Inches(7.7), Inches(3.35), STEEL, 2)
    add_rich_lines(slide, [
        ("RESULTADO  ", "ejecución inmutable, reintentos seguros y revisión accionable", OXIDE),
        ("CAPACIDAD  ", "2 desarrolladores × 6 días enfocados", CHALK_MUTED),
        ("ALTERNATIVA  ", "1 desarrollador × 9–10 días enfocados", CHALK_MUTED),
    ], Inches(0.82), Inches(3.72), Inches(7.1), Inches(1.7), size=16, gap=9)
    add_rect(slide, Inches(9.0), Inches(1.05), Inches(3.3), Inches(4.9), GRAPHITE_2, STEEL)
    for i, (name, color) in enumerate([
        ("PROGRAMAR", OXIDE), ("EJECUTAR", CHALK_MUTED), ("DETECTAR", WARNING),
        ("REVISAR", CHALK_MUTED), ("AJUSTAR", SUCCESS),
    ]):
        y = Inches(1.45 + i * 0.86)
        add_text(slide, f"0{i + 1}", Inches(9.35), y, Inches(0.45), Inches(0.28), size=10, color=STEEL_LIGHT, bold=True)
        add_text(slide, name, Inches(9.95), y - Inches(0.03), Inches(1.6), Inches(0.32), size=16, color=color, bold=True)
        if i < 4:
            add_line(slide, Inches(9.57), y + Inches(0.28), Inches(9.57), y + Inches(0.72), STEEL, 1)
    add_text(slide, "Guía para desarrollo · MVP piloto", Inches(0.8), Inches(7.08), Inches(3.5), Inches(0.2), size=8, color=STEEL_LIGHT)
    add_text(slide, "01", Inches(12.1), Inches(7.05), Inches(0.5), Inches(0.2), size=9, color=CHALK_MUTED, bold=True, align=PP_ALIGN.RIGHT)

    # 2
    slide = base_slide(prs, "Tesis del producto", "Lo que debe demostrar el piloto", 2)
    add_text(slide, "Forja es el ciclo operativo entre la prescripción y el ajuste.",
             Inches(0.7), Inches(1.58), Inches(11.6), Inches(0.55), size=25, bold=True)
    stages = [("PROGRAMAR", "Intención del coach"), ("EJECUTAR", "Acción del atleta"),
              ("DETECTAR", "Desvío relevante"), ("REVISAR", "Criterio del coach"),
              ("AJUSTAR", "Siguiente prescripción")]
    for i, (name, body) in enumerate(stages):
        x = Inches(0.72 + i * 2.45)
        add_rect(slide, x, Inches(2.75), Inches(1.92), Inches(1.25), GRAPHITE_2, STEEL, radius=True)
        add_text(slide, f"0{i+1}", x + Inches(0.18), Inches(2.94), Inches(0.35), Inches(0.2), size=9, color=OXIDE, bold=True)
        add_text(slide, name, x + Inches(0.18), Inches(3.2), Inches(1.5), Inches(0.25), size=14, bold=True)
        add_text(slide, body, x + Inches(0.18), Inches(3.53), Inches(1.55), Inches(0.25), size=10, color=CHALK_MUTED)
        if i < 4:
            add_line(slide, x + Inches(1.92), Inches(3.36), x + Inches(2.42), Inches(3.36), OXIDE, 2)
    label_box(slide, "Prueba piloto", "Una sesión asignada se convierte en evidencia durable, revisión priorizada y ajuste deliberado.",
              Inches(1.2), Inches(4.75), Inches(10.9), Inches(1.2), accent=SUCCESS, size=17)

    # 3
    slide = base_slide(prs, "Tres actores, un solo ciclo", "Modelo del producto", 3)
    actor_data = [
        ("DUEÑO / COMPRADOR", "Comprende el estado operativo\ny confía en el uso continuo", OXIDE),
        ("COACH OPERADOR", "Prescribe, detecta, revisa\ny ajusta", WARNING),
        ("ATLETA", "Ejecuta y registra\ncon fricción mínima", SUCCESS),
    ]
    for i, (title, body, color) in enumerate(actor_data):
        x = Inches(0.75 + i * 4.15)
        label_box(slide, title, body, x, Inches(1.72), Inches(3.55), Inches(1.55), color, 16)
    add_line(slide, Inches(2.5), Inches(3.55), Inches(10.75), Inches(3.55), STEEL, 2)
    add_text(slide, "comprender", Inches(1.6), Inches(3.35), Inches(1.2), Inches(0.25), size=9, color=CHALK_MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "prescribir →", Inches(5.35), Inches(3.35), Inches(1.4), Inches(0.25), size=9, color=CHALK_MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "← evidencia", Inches(9.5), Inches(3.35), Inches(1.4), Inches(0.25), size=9, color=CHALK_MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Regla de optimización", Inches(0.75), Inches(4.35), Inches(2.5), Inches(0.3), size=12, color=OXIDE, bold=True)
    add_text(slide, "Nunca mejorar la experiencia de un actor debilitando la confianza de otro en los mismos datos.",
             Inches(0.75), Inches(4.78), Inches(11.5), Inches(0.6), size=23, bold=True)
    add_text(slide, "Los equipos importan cuando son necesarios. La multitenencia empresarial no demuestra este piloto.",
             Inches(0.75), Inches(5.72), Inches(10.5), Inches(0.35), size=14, color=CHALK_MUTED)

    # 4
    slide = base_slide(prs, "Recorridos y capacidades actuales", "Evidencia, no aspiración", 4)
    cols = [
        ("COACH", ["Inicio de sesión + autorización", "Panel + atletas", "Planificación + marcas", "Historial + E/S de Excel"]),
        ("ATLETA", ["Activación PIN / OTP", "Entrenamiento de hoy", "Registro de series + día", "Vista de progreso"]),
        ("SISTEMA", ["Drizzle / libSQL", "Auth.js", "Evolution API", "Cola offline parcial"]),
    ]
    for i, (name, items) in enumerate(cols):
        x = Inches(0.72 + i * 4.16)
        add_text(slide, name, x, Inches(1.55), Inches(3.4), Inches(0.3), size=11, color=OXIDE, bold=True)
        for j, item in enumerate(items):
            y = Inches(2.05 + j * 0.72)
            add_rect(slide, x, y, Inches(3.55), Inches(0.52), GRAPHITE_2, STEEL, radius=True)
            add_text(slide, item, x + Inches(0.18), y + Inches(0.11), Inches(3.15), Inches(0.25), size=13, color=CHALK)
    add_rect(slide, Inches(0.72), Inches(5.3), Inches(11.9), Inches(0.9), OXIDE_DARK, OXIDE_DARK)
    add_text(slide, "La base es suficiente. El límite del piloto es la integridad y el foco operativo, no la cantidad de funciones.",
             Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.35), size=17, bold=True, align=PP_ALIGN.CENTER)

    # 5
    slide = base_slide(prs, "Arquitectura actual del sistema", "Dónde aparece el acoplamiento", 5)
    label_box(slide, "RUTAS / ACCIONES", "Superficies App Router de coach + atleta", Inches(0.75), Inches(1.65), Inches(3.1), Inches(1.05), OXIDE, 14)
    label_box(slide, "GRAFO DEL PLAN", "programa → semana → día → ejercicio → serie", Inches(5.1), Inches(1.65), Inches(3.1), Inches(1.05), WARNING, 14)
    label_box(slide, "LECTURAS HISTÓRICAS", "registros unidos mediante un plan mutable", Inches(9.45), Inches(1.65), Inches(3.1), Inches(1.05), OXIDE, 14)
    add_line(slide, Inches(3.85), Inches(2.18), Inches(5.1), Inches(2.18), STEEL_LIGHT, 2)
    add_line(slide, Inches(8.2), Inches(2.18), Inches(9.45), Inches(2.18), STEEL_LIGHT, 2)
    add_text(slide, "RIESGO DE CASCADA / MUTACIÓN", Inches(4.35), Inches(3.14), Inches(4.65), Inches(0.3), size=12, color=OXIDE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.25), Inches(3.62), Inches(6.85), Inches(1.1), GRAPHITE_2, OXIDE)
    add_text(slide, "El rendimiento completado no tiene una instantánea de sesión independiente.",
             Inches(3.55), Inches(3.91), Inches(6.25), Inches(0.35), size=20, bold=True, align=PP_ALIGN.CENTER)
    add_rich_lines(slide, [
        ("ADEMÁS  ", "las transacciones y restricciones están incompletas", WARNING),
        ("ADEMÁS  ", "los reintentos offline carecen de un contrato idempotente", WARNING),
        ("ADEMÁS  ", "las semánticas de calendario y pendientes divergen", WARNING),
    ], Inches(3.65), Inches(5.1), Inches(6.0), Inches(1.3), size=13, gap=5)

    # 6
    slide = base_slide(prs, "Bases que debemos conservar", "Sin reescritura", 6)
    foundations = [
        ("STACK", "Next.js 16 · TypeScript · Tailwind v4", OXIDE),
        ("DATOS", "Drizzle + libSQL", WARNING),
        ("IDENTIDAD", "Auth.js + propiedad verificada en servidor", SUCCESS),
        ("ACTIVACIÓN", "OTP de WhatsApp vía Evolution API", OXIDE),
        ("OPERACIONES", "Puente Excel + cola offline parcial", WARNING),
        ("DOMINIO", "Jerarquía del programa + marcas de fuerza", SUCCESS),
    ]
    for i, (title, body, color) in enumerate(foundations):
        col, row = i % 2, i // 2
        x = Inches(0.78 + col * 6.15)
        y = Inches(1.6 + row * 1.45)
        label_box(slide, title, body, x, y, Inches(5.55), Inches(1.05), color, 15)
    add_text(slide, "Dirección: monolito modular con operaciones explícitas de dominio, no nuevas unidades de despliegue.",
             Inches(0.8), Inches(6.05), Inches(11.4), Inches(0.35), size=17, color=CHALK, bold=True, align=PP_ALIGN.CENTER)

    # 7
    slide = base_slide(prs, "Brechas críticas", "Bloqueantes del piloto", 7)
    gaps = [
        ("01", "Plan mutable ↔ historial", "Pérdida de datos / evidencia reescrita"),
        ("02", "Sin instantáneas de sesión", "No hay agregado durable de ejecución"),
        ("03", "Garantías débiles de escritura", "Riesgo de reintento, parcialidad y duplicación"),
        ("04", "El panel no es una bandeja", "Conteos sin siguiente acción"),
        ("05", "%RM sin resolver", "Prescripción ambigua al ejecutar"),
        ("06", "Evidencia de lanzamiento incompleta", "Brechas de CI, E2E, restauración y operación"),
    ]
    for i, (num, title, impact) in enumerate(gaps):
        col, row = i % 2, i // 2
        x = Inches(0.75 + col * 6.2)
        y = Inches(1.55 + row * 1.48)
        add_text(slide, num, x, y + Inches(0.08), Inches(0.55), Inches(0.3), size=12, color=OXIDE, bold=True)
        add_text(slide, title, x + Inches(0.7), y, Inches(4.8), Inches(0.32), size=17, bold=True)
        add_text(slide, impact, x + Inches(0.7), y + Inches(0.43), Inches(4.8), Inches(0.28), size=12, color=CHALK_MUTED)
        add_line(slide, x + Inches(0.7), y + Inches(0.86), x + Inches(5.45), y + Inches(0.86), STEEL, 1)
    add_rect(slide, Inches(0.75), Inches(6.02), Inches(11.8), Inches(0.48), OXIDE_DARK, OXIDE_DARK)
    add_text(slide, "Detención obligatoria: no hay piloto si historial, idempotencia, autorización, migración o restauración están en rojo.",
             Inches(0.95), Inches(6.11), Inches(11.4), Inches(0.25), size=13, bold=True, align=PP_ALIGN.CENTER)

    # 8
    slide = base_slide(prs, "Límite del MVP piloto", "Lo obligatorio prevalece", 8)
    scope = [
        ("DEBE", SUCCESS, ["Instantáneas de sesión + rendimiento", "Estados + transiciones seguras", "Transacciones + idempotencia", "Bandeja + calendario canónico", "Instantánea %RM + logout", "Migración / restauración / E2E"]),
        ("DEBERÍA", WARNING, ["Duplicar unidades del plan", "Confirmación de revisión", "Estado del sistema", "Edición de series más rápida"]),
        ("DESPUÉS", STEEL_LIGHT, ["Multitenencia empresarial", "IA / chat / nutrición", "Facturación / marketplace", "Rebranding final"]),
    ]
    widths = [4.55, 3.55, 3.35]
    x = 0.72
    for idx, ((title, color, items), width) in enumerate(zip(scope, widths)):
        add_text(slide, title, Inches(x), Inches(1.55), Inches(width), Inches(0.3), size=12, color=color, bold=True)
        add_rect(slide, Inches(x), Inches(1.98), Inches(width), Inches(4.2), GRAPHITE_2, color)
        for j, item in enumerate(items):
            add_text(slide, "—", Inches(x + 0.23), Inches(2.25 + j * 0.55), Inches(0.25), Inches(0.25), size=12, color=color, bold=True)
            add_text(slide, item, Inches(x + 0.55), Inches(2.24 + j * 0.55), Inches(width - 0.75), Inches(0.38), size=12 if idx == 0 else 11, color=CHALK)
        x += width + 0.35

    # 9
    slide = base_slide(prs, "Ejecución paralela en seis días", "Dos frentes, integración diaria", 9)
    headers = [("DÍA", 0.75, 0.7), ("DESARROLLADOR A", 1.6, 4.85), ("DESARROLLADOR B", 6.65, 4.85), ("DEMO", 11.65, 0.85)]
    for text, x, w in headers:
        add_text(slide, text, Inches(x), Inches(1.42), Inches(w), Inches(0.25), size=9, color=OXIDE, bold=True)
    rows = [
        ("1", "Esquema + ciclo de vida", "Semántica + fixtures", "contrato"),
        ("2", "Comandos transaccionales", "Atleta + adaptador offline", "guardado"),
        ("3", "Lecturas de instantáneas", "Resolución %RM", "historial"),
        ("4", "Bandeja del coach", "Alineación de calendario", "revisión"),
        ("5", "Revisión + planificación", "Logout + robustez de sync", "recorrido"),
        ("6", "Migración + recuperación", "E2E + accesibilidad", "evidencia"),
    ]
    for i, (day, a, b, demo) in enumerate(rows):
        y = Inches(1.82 + i * 0.76)
        fill = GRAPHITE_2 if i % 2 == 0 else GRAPHITE
        add_rect(slide, Inches(0.7), y, Inches(11.95), Inches(0.58), fill, STEEL)
        add_text(slide, day, Inches(0.87), y + Inches(0.11), Inches(0.3), Inches(0.25), size=13, color=OXIDE, bold=True)
        add_text(slide, a, Inches(1.6), y + Inches(0.11), Inches(4.75), Inches(0.25), size=12, color=CHALK)
        add_text(slide, b, Inches(6.65), y + Inches(0.11), Inches(4.7), Inches(0.25), size=12, color=CHALK)
        add_text(slide, demo, Inches(11.55), y + Inches(0.11), Inches(0.9), Inches(0.25), size=10, color=SUCCESS, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Alternativa con una persona: 9–10 días · mismo orden · quitar primero lo opcional.",
             Inches(0.75), Inches(6.55), Inches(11.5), Inches(0.28), size=13, color=CHALK_MUTED, bold=True, align=PP_ALIGN.CENTER)

    # 10
    slide = base_slide(prs, "Dominio y modelo de datos objetivo", "Ejecución inmutable", 10)
    label_box(slide, "PROGRAMA", "draft → published → archived", Inches(0.75), Inches(1.55), Inches(3.2), Inches(1.05), OXIDE, 15)
    label_box(slide, "VERSIÓN DEL PLAN", "día / ejercicio / serie planificados", Inches(0.75), Inches(3.08), Inches(3.2), Inches(1.05), WARNING, 15)
    add_line(slide, Inches(2.35), Inches(2.6), Inches(2.35), Inches(3.08), STEEL_LIGHT, 2)
    add_line(slide, Inches(3.95), Inches(3.6), Inches(5.0), Inches(3.6), OXIDE, 2)
    label_box(slide, "workout_sessions", "scheduled → in_progress → completed → reviewed", Inches(5.0), Inches(2.35), Inches(3.35), Inches(1.55), SUCCESS, 14)
    add_line(slide, Inches(8.35), Inches(3.12), Inches(9.4), Inches(3.12), OXIDE, 2)
    label_box(slide, "set_performances", "instantánea de prescripción\ncarga resuelta + ejecución\ndesvíos + clave de operación", Inches(9.4), Inches(1.9), Inches(3.15), Inches(2.45), OXIDE, 14)
    add_rich_lines(slide, [
        ("INVARIANTE  ", "editar el plan nunca reescribe evidencia completada", OXIDE),
        ("REINTENTO  ", "la clave estable devuelve el resultado aceptado", WARNING),
        ("COMMIT  ", "rendimientos + cierre + evento de bandeja son atómicos", SUCCESS),
    ], Inches(1.1), Inches(5.0), Inches(11.2), Inches(1.25), size=14, gap=5)

    # 11
    slide = base_slide(prs, "Verificación y condiciones de salida", "Evidencia antes del piloto", 11)
    gates = [
        ("DATOS", "Migración + restricciones + restauración", OXIDE),
        ("COMANDO", "Duplicación + fallo parcial", WARNING),
        ("HISTORIAL", "Cambiar el plan no altera la instantánea", SUCCESS),
        ("ACCESO", "Rechazo cruzado + sesión cerrada", OXIDE),
        ("RECORRIDO", "Publicar → ejecutar → revisar → ajustar", WARNING),
        ("OPERACIÓN", "Build + smoke + ensayo de rollback", SUCCESS),
    ]
    for i, (name, proof, color) in enumerate(gates):
        col, row = i % 2, i // 2
        x = Inches(0.75 + col * 6.18)
        y = Inches(1.55 + row * 1.35)
        add_rect(slide, x, y, Inches(5.55), Inches(0.98), GRAPHITE_2, STEEL, radius=True)
        add_rect(slide, x + Inches(0.18), y + Inches(0.22), Inches(0.52), Inches(0.52), color, color, radius=True)
        add_text(slide, "✓", x + Inches(0.18), y + Inches(0.23), Inches(0.52), Inches(0.4), size=16, color=GRAPHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, name, x + Inches(0.9), y + Inches(0.16), Inches(1.2), Inches(0.25), size=10, color=color, bold=True)
        add_text(slide, proof, x + Inches(0.9), y + Inches(0.48), Inches(4.25), Inches(0.25), size=13, color=CHALK)
    add_rect(slide, Inches(0.75), Inches(5.8), Inches(11.75), Inches(0.62), OXIDE_DARK, OXIDE_DARK)
    add_text(slide, "La decisión es binaria: todo lo obligatorio en verde o reducir alcance y repetir la evidencia.",
             Inches(1.0), Inches(5.96), Inches(11.25), Inches(0.28), size=16, bold=True, align=PP_ALIGN.CENTER)

    # 12
    slide = base_slide(prs, "Traspaso y punto de partida recomendado", "Primer movimiento", 12)
    add_text(slide, "Empezar por el contrato de ejecución, no por el panel.",
             Inches(0.75), Inches(1.53), Inches(11.8), Inches(0.55), size=27, bold=True)
    steps = [
        ("01", "Fijar el fixture piloto", "actor, calendario, %RM, desviación y reintento"),
        ("02", "Agregar instantáneas", "sesión + rendimiento + restricciones"),
        ("03", "Controlar los comandos", "autorización + transacción + idempotencia"),
        ("04", "Migrar las proyecciones", "primero historial, luego bandeja"),
        ("05", "Demostrar el lanzamiento", "migración + recuperación + recorrido completo"),
    ]
    for i, (num, title, detail) in enumerate(steps):
        y = Inches(2.35 + i * 0.72)
        add_text(slide, num, Inches(0.82), y, Inches(0.5), Inches(0.3), size=11, color=OXIDE, bold=True)
        add_text(slide, title, Inches(1.55), y - Inches(0.02), Inches(3.45), Inches(0.32), size=16, bold=True)
        add_text(slide, detail, Inches(5.2), y, Inches(5.9), Inches(0.28), size=13, color=CHALK_MUTED)
        add_line(slide, Inches(1.55), y + Inches(0.42), Inches(12.0), y + Inches(0.42), STEEL, 1)
    add_rect(slide, Inches(8.85), Inches(5.95), Inches(3.35), Inches(0.62), SUCCESS, SUCCESS)
    add_text(slide, "SIGUIENTE: CONTRATO DEL DÍA 1", Inches(8.85), Inches(6.1), Inches(3.35), Inches(0.28), size=13, color=GRAPHITE, bold=True, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(f"Generado: {path}")
